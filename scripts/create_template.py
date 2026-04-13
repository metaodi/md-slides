"""Script to generate the default md-slides template.pptx.

Run this script to regenerate the bundled template:

    python scripts/create_template.py

The output is written to md_slides/template.pptx so that it is included as
package data when the project is installed.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# Output path inside the package so it ships as package data
OUTPUT = Path(__file__).parent.parent / "md_slides" / "template.pptx"

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1F, 0x49, 0x7D)  # title-slide background
ACCENT_BLUE = RGBColor(0x27, 0x6F, 0xBF)  # title-slide subtitle / accents
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY = RGBColor(0x26, 0x26, 0x26)  # body text on content slides
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)  # content-slide background


def _set_fill_solid(fill_element, rgb: RGBColor):
    """Set a solid colour fill on an XML fill element."""
    # Remove any existing fill children
    for child in list(fill_element):
        fill_element.remove(child)
    solidFill = fill_element.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": f"{rgb:06X}"})
    solidFill.append(srgbClr)
    fill_element.append(solidFill)


def _set_layout_background(layout, rgb: RGBColor):
    """Set the background colour of a slide layout."""
    bg = layout.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _set_text_colour(text_frame, rgb: RGBColor, font_size_pt: int = None):
    """Set colour (and optionally size) for all runs in a text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = rgb
            if font_size_pt is not None:
                run.font.size = Pt(font_size_pt)


def build_template(output_path: Path):
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_master = prs.slide_master

    # ── Layout 0: Title Slide ─────────────────────────────────────────────────
    title_layout = slide_master.slide_layouts[0]
    _set_layout_background(title_layout, DARK_BLUE)

    # Style every placeholder on the title layout
    for ph in title_layout.placeholders:
        idx = ph.placeholder_format.idx
        tf = ph.text_frame
        if idx == 0:
            # Main title placeholder
            for para in tf.paragraphs:
                para.runs  # ensure paragraph exists
                pPr = para._p.get_or_add_pPr()
                for run in para.runs:
                    run.font.color.rgb = WHITE
                    run.font.size = Pt(44)
                    run.font.bold = True
        elif idx == 1:
            # Subtitle placeholder
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xBF, 0xD7, 0xFF)
                    run.font.size = Pt(24)

    # ── Layout 1: Title and Content ───────────────────────────────────────────
    content_layout = slide_master.slide_layouts[1]
    _set_layout_background(content_layout, WHITE)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Template saved to {output_path}")


if __name__ == "__main__":
    build_template(OUTPUT)
    sys.exit(0)
