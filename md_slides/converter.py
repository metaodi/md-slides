"""PPTX converter for md-slides.

Converts a list of slide dicts (produced by the parser) into a PowerPoint
presentation using python-pptx and a template file.
"""

import importlib.resources
import os

from pptx import Presentation
from pptx.util import Pt

# Indices into the template's slide layout collection.
# Layout 0: Title Slide  (used for slides with type='title')
# Layout 1: Title and Content  (used for slides with type='content')
TITLE_SLIDE_LAYOUT_INDEX = 0
CONTENT_SLIDE_LAYOUT_INDEX = 1


def get_default_template_path():
    """Return the absolute path to the bundled default template.

    Returns:
        Absolute path string to the packaged template.pptx.
    """
    # Python 3.9+: use importlib.resources.files
    try:
        ref = importlib.resources.files("md_slides").joinpath("template.pptx")
        with importlib.resources.as_file(ref) as path:
            return str(path)
    except (AttributeError, TypeError):
        # Fallback for older Python
        import pkg_resources

        return pkg_resources.resource_filename("md_slides", "template.pptx")


def convert(slides, output_path, template_path=None):
    """Convert a list of slide dicts to a PowerPoint file.

    Args:
        slides: List of slide dicts from the parser.
        output_path: Destination .pptx file path.
        template_path: Path to a .pptx template. Uses bundled default if None.
    """
    if template_path is None:
        template_path = get_default_template_path()

    prs = Presentation(template_path)

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        if slide_type == "title":
            _add_title_slide(prs, slide_data)
        else:
            _add_content_slide(prs, slide_data)
    
    if not output_path.endswith(".pptx"):
        output_path = f"{output_path}.pptx"

    prs.save(output_path)
    return output_path


def _add_title_slide(prs, slide_data):
    """Add a title slide to the presentation.

    Args:
        prs: python-pptx Presentation object.
        slide_data: Dict with 'title' and optional 'subtitle' keys.
    """
    layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)

    title_ph = _get_placeholder(slide, 0)
    if title_ph is not None:
        title_ph.text = slide_data.get("title", "")

    subtitle_ph = _get_placeholder(slide, 1)
    if subtitle_ph is not None:
        subtitle = slide_data.get("subtitle", "")
        subtitle_ph.text = subtitle


def _add_content_slide(prs, slide_data):
    """Add a content slide to the presentation.

    Args:
        prs: python-pptx Presentation object.
        slide_data: Dict with 'title' and 'elements' keys.
    """
    layout = prs.slide_layouts[CONTENT_SLIDE_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)

    title_ph = _get_placeholder(slide, 0)
    if title_ph is not None:
        title_ph.text = slide_data.get("title", "")

    body_ph = _get_placeholder(slide, 1)
    if body_ph is not None:
        _fill_body(body_ph.text_frame, slide_data.get("elements", []))


def _fill_body(text_frame, elements):
    """Fill a text frame with parsed body elements.

    Args:
        text_frame: python-pptx TextFrame object.
        elements: List of element dicts from the parser.
    """
    # Clear the default empty paragraph
    text_frame.clear()

    first = True
    for elem in elements:
        if first:
            # Re-use the first (cleared) paragraph
            para = text_frame.paragraphs[0]
            first = False
        else:
            para = text_frame.add_paragraph()

        if elem["type"] == "bullet":
            para.level = elem["level"] 

        _apply_runs(para, elem.get("runs", []))


def _apply_runs(paragraph, runs):
    """Apply formatted runs to a paragraph.

    Args:
        paragraph: python-pptx Paragraph object.
        runs: List of run dicts with 'text', 'bold', 'italic'.
    """
    for run_data in runs:
        run = paragraph.add_run()
        run.text = run_data.get("text", "")
        if run_data.get("bold"):
            run.font.bold = True
        if run_data.get("italic"):
            run.font.italic = True


def _get_placeholder(slide, idx):
    """Return a placeholder by index, or None if not present.

    Args:
        slide: python-pptx Slide object.
        idx: Placeholder index (0=title, 1=body/subtitle).

    Returns:
        Placeholder shape or None.
    """
    placeholders = list(slide.placeholders)
    return placeholders[idx] if idx < len(placeholders) else None
