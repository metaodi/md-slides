"""Tests for md_slides.converter."""

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from md_slides.converter import convert, get_default_template_path
from md_slides.parser import parse_markdown

# ── Helper ─────────────────────────────────────────────────────────────────────


def _make_pptx(md_content):
    """Parse *md_content* and convert it to a temporary PPTX; return the path."""
    slides = parse_markdown(md_content)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fh:
        out = fh.name
    convert(slides, out)
    return out


# ── Template ───────────────────────────────────────────────────────────────────


def test_default_template_exists():
    path = get_default_template_path()
    assert os.path.isfile(path), f"Default template not found at {path}"


# ── Slide count ────────────────────────────────────────────────────────────────


def test_single_title_slide():
    out = _make_pptx("# Hello World")
    try:
        prs = Presentation(out)
        assert len(prs.slides) == 1
    finally:
        os.unlink(out)


def test_multiple_slides_correct_count():
    md = "# Title\nSubtitle\n\n---\n\n## Content\n- bullet\n\n---\n\n## Second\n- item"
    out = _make_pptx(md)
    try:
        prs = Presentation(out)
        assert len(prs.slides) == 3
    finally:
        os.unlink(out)


# ── Layout selection ───────────────────────────────────────────────────────────


def test_title_slide_uses_title_layout():
    """Title slide should use layout index 0 ('Title Slide')."""
    out = _make_pptx("# My Title\nMy subtitle")
    try:
        prs = Presentation(out)
        slide = prs.slides[0]
        assert slide.slide_layout.name == "Title Slide"
    finally:
        os.unlink(out)


def test_content_slide_uses_content_layout():
    """Content slide should use layout index 1 ('Title and Content')."""
    out = _make_pptx("## Content Slide\n- bullet")
    try:
        prs = Presentation(out)
        slide = prs.slides[0]
        assert slide.slide_layout.name == "Title and Content"
    finally:
        os.unlink(out)


# ── Content ────────────────────────────────────────────────────────────────────


def test_title_slide_title_text():
    out = _make_pptx("# My Title\nMy subtitle")
    try:
        prs = Presentation(out)
        slide = prs.slides[0]
        title_ph = next(
            ph for ph in slide.placeholders if ph.placeholder_format.idx == 0
        )
        assert title_ph.text == "My Title"
    finally:
        os.unlink(out)


def test_title_slide_subtitle_text():
    out = _make_pptx("# My Title\nMy subtitle")
    try:
        prs = Presentation(out)
        slide = prs.slides[0]
        subtitle_ph = next(
            ph for ph in slide.placeholders if ph.placeholder_format.idx == 1
        )
        assert subtitle_ph.text == "My subtitle"
    finally:
        os.unlink(out)


def test_content_slide_title_text():
    out = _make_pptx("## My Content Slide\n- item")
    try:
        prs = Presentation(out)
        slide = prs.slides[0]
        title_ph = next(
            ph for ph in slide.placeholders if ph.placeholder_format.idx == 0
        )
        assert title_ph.text == "My Content Slide"
    finally:
        os.unlink(out)


def test_content_slide_body_contains_bullets():
    out = _make_pptx("## Slide\n- First\n- Second\n- Third")
    try:
        prs = Presentation(out)
        slide = prs.slides[0]
        body_ph = next(
            ph for ph in slide.placeholders if ph.placeholder_format.idx == 1
        )
        full_text = body_ph.text_frame.text
        assert "First" in full_text
        assert "Second" in full_text
        assert "Third" in full_text
    finally:
        os.unlink(out)


def test_bold_text_in_bullet():
    out = _make_pptx("## Slide\n- **bold item**")
    try:
        prs = Presentation(out)
        slide = prs.slides[0]
        body_ph = next(
            ph for ph in slide.placeholders if ph.placeholder_format.idx == 1
        )
        # Find a bold run
        bold_runs = [
            run
            for para in body_ph.text_frame.paragraphs
            for run in para.runs
            if run.font.bold
        ]
        assert any(r.text == "bold item" for r in bold_runs)
    finally:
        os.unlink(out)


def test_italic_text_in_bullet():
    out = _make_pptx("## Slide\n- *italic item*")
    try:
        prs = Presentation(out)
        slide = prs.slides[0]
        body_ph = next(
            ph for ph in slide.placeholders if ph.placeholder_format.idx == 1
        )
        italic_runs = [
            run
            for para in body_ph.text_frame.paragraphs
            for run in para.runs
            if run.font.italic
        ]
        assert any(r.text == "italic item" for r in italic_runs)
    finally:
        os.unlink(out)


# ── Custom template ────────────────────────────────────────────────────────────


def test_custom_template(tmp_path):
    """Convert with an explicitly supplied template path."""
    template_path = get_default_template_path()
    out = str(tmp_path / "out.pptx")
    slides = parse_markdown("# Title\nSubtitle\n\n---\n\n## Content\n- bullet")
    convert(slides, out, template_path=template_path)
    prs = Presentation(out)
    assert len(prs.slides) == 2
